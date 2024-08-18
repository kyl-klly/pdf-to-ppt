import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_png(pdf_path, output_folder, dpi=300):
    # Open the PDF file
    pdf_document = fitz.open(pdf_path)
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # Convert each page to a PNG image with higher resolution
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        mat = fitz.Matrix(dpi / 72, dpi / 72)  # Set the resolution
        pix = page.get_pixmap(matrix=mat)
        output_image_path = os.path.join(output_folder, f"page_{page_num + 1:04d}.png")
        pix.save(output_image_path)
        print(f"Saved {output_image_path}")

def create_ppt_from_images(image_folder, output_ppt):
    # Create a presentation object
    presentation = Presentation()
    
    # Add each image as a slide
    for image_file in sorted(os.listdir(image_folder)):
        if image_file.endswith(".png"):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(os.path.join(image_folder, image_file), left, top, width=presentation.slide_width, height=presentation.slide_height)
    
    # Save the presentation
    presentation.save(output_ppt)
    print(f"Presentation saved as {output_ppt}")

if __name__ == "__main__":
    pdf_path = "PageSizes_output.pdf"  # Replace with your PDF file path
    output_folder = "output_images"
    output_ppt = "ppt_file.pptx"
    
    pdf_to_png(pdf_path, output_folder)
    create_ppt_from_images(output_folder, output_ppt)
